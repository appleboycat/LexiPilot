#!/usr/bin/env python3
"""Build LexiPilot submission PDF/PPTX artifacts from committed sources."""

from __future__ import annotations

import argparse
import html
import os
import re
import shutil
import tempfile
from pathlib import Path
from typing import Any, Iterable

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER, TA_LEFT
from reportlab.lib.pagesizes import A4, landscape
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.platypus import (
    Image as RLImage,
    KeepTogether,
    PageBreak,
    Paragraph,
    Preformatted,
    SimpleDocTemplate,
    Spacer,
    Table,
    TableStyle,
)


ROOT = Path(__file__).resolve().parents[1]
SUBMISSION = ROOT / "submission"
ARCH_PNG = SUBMISSION / "architecture" / "lexipilot_architecture.png"
SPEC_MD = SUBMISSION / "LexiPilot_Project_Specification.md"
SPEC_PDF = SUBMISSION / "LexiPilot_Project_Specification.pdf"
PRESENTATION = SUBMISSION / "LexiPilot_Presentation.pptx"
PRESENTATION_PDF = SUBMISSION / "LexiPilot_Presentation.pdf"
OFFICIAL = ROOT / "official_submission" / "LexiPilot"

WIDTH, HEIGHT = 1920, 1080
BG = "#11161B"
PANEL = "#1B232B"
PANEL_ALT = "#202B34"
WHITE = "#F5F7F8"
MUTED = "#AEB8C2"
RED = "#E13B2C"
CYAN = "#43B8CF"
GREEN = "#64C878"
YELLOW = "#F3C74F"
PURPLE = "#B99BE7"
LINE = "#34424E"

FONT_REGULAR = Path("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf")
FONT_BOLD = Path("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf")
FONT_MONO = Path("/usr/share/fonts/truetype/dejavu/DejaVuSansMono.ttf")


def font(size: int, *, bold: bool = False, mono: bool = False) -> ImageFont.FreeTypeFont:
    path = FONT_MONO if mono else FONT_BOLD if bold else FONT_REGULAR
    return ImageFont.truetype(str(path), size)


def atomic_replace(source: Path, destination: Path) -> None:
    destination.parent.mkdir(parents=True, exist_ok=True)
    os.replace(source, destination)


def wrap_text(draw: ImageDraw.ImageDraw, text: str, selected_font: ImageFont.FreeTypeFont, width: int) -> list[str]:
    words = text.split()
    lines: list[str] = []
    current = ""
    for word in words:
        candidate = f"{current} {word}".strip()
        if not current or draw.textbbox((0, 0), candidate, font=selected_font)[2] <= width:
            current = candidate
        else:
            lines.append(current)
            current = word
    if current:
        lines.append(current)
    return lines


def draw_wrapped(
    draw: ImageDraw.ImageDraw,
    xy: tuple[int, int],
    text: str,
    selected_font: ImageFont.FreeTypeFont,
    fill: str,
    width: int,
    *,
    line_gap: int = 10,
) -> int:
    x, y = xy
    line_height = selected_font.size + line_gap
    for line in wrap_text(draw, text, selected_font, width):
        draw.text((x, y), line, font=selected_font, fill=fill)
        y += line_height
    return y


def slide_base(number: int, title: str, kicker: str = "") -> tuple[Image.Image, ImageDraw.ImageDraw]:
    image = Image.new("RGB", (WIDTH, HEIGHT), BG)
    draw = ImageDraw.Draw(image)
    draw.rectangle((0, 0, WIDTH, 14), fill=RED)
    if kicker:
        draw.text((90, 62), kicker.upper(), font=font(24, bold=True), fill=CYAN)
    draw.text((90, 105), title, font=font(54, bold=True), fill=WHITE)
    draw.line((90, 182, 1830, 182), fill=LINE, width=2)
    draw.text((90, 1025), "LexiPilot | AMD Radeon Hackathon 2026 | Track 2", font=font(20), fill=MUTED)
    draw.text((1780, 1025), f"{number}/8", font=font(20, bold=True), fill=MUTED)
    return image, draw


def panel(draw: ImageDraw.ImageDraw, box: tuple[int, int, int, int], *, outline: str = LINE, fill: str = PANEL) -> None:
    draw.rounded_rectangle(box, radius=10, fill=fill, outline=outline, width=2)


def bullet_list(
    draw: ImageDraw.ImageDraw,
    items: Iterable[str],
    box: tuple[int, int, int, int],
    *,
    accent: str = CYAN,
    text_size: int = 30,
    gap: int = 25,
) -> None:
    x1, y1, x2, _ = box
    y = y1
    selected_font = font(text_size)
    for item in items:
        draw.ellipse((x1, y + 11, x1 + 12, y + 23), fill=accent)
        y = draw_wrapped(draw, (x1 + 30, y), item, selected_font, WHITE, x2 - x1 - 35, line_gap=8) + gap


def card(
    draw: ImageDraw.ImageDraw,
    box: tuple[int, int, int, int],
    heading: str,
    body: str,
    accent: str,
) -> None:
    panel(draw, box)
    x1, y1, x2, _ = box
    draw.rectangle((x1, y1, x1 + 8, box[3]), fill=accent)
    draw.text((x1 + 35, y1 + 28), heading, font=font(28, bold=True), fill=accent)
    draw_wrapped(draw, (x1 + 35, y1 + 82), body, font(25), WHITE, x2 - x1 - 70, line_gap=7)


def slide_1() -> Image.Image:
    image = Image.new("RGB", (WIDTH, HEIGHT), BG)
    draw = ImageDraw.Draw(image)
    draw.rectangle((0, 0, 18, HEIGHT), fill=RED)
    draw.text((115, 120), "LEXIPILOT", font=font(86, bold=True), fill=WHITE)
    draw.text((118, 235), "Private Adaptive Vocabulary Learning Agent", font=font(43), fill=CYAN)
    draw.text((118, 305), "on AMD Radeon", font=font(43, bold=True), fill=RED)
    draw.line((118, 390, 1130, 390), fill=LINE, width=3)
    draw_wrapped(
        draw,
        (118, 440),
        "Model-selected read-only tools. Validated structured plans. Deterministic progress writes.",
        font(36),
        WHITE,
        1040,
        line_gap=12,
    )
    panel(draw, (1270, 120, 1815, 860), outline=RED, fill=PANEL)
    draw.text((1330, 180), "TRACK 2", font=font(31, bold=True), fill=RED)
    bullet_list(
        draw,
        [
            "Qwen3-8B planning",
            "vLLM + ROCm",
            "Dedicated AMD Radeon inference",
            "Persistent learner memory",
            "Safe deterministic controller",
        ],
        (1330, 260, 1755, 810),
        accent=CYAN,
        text_size=27,
        gap=24,
    )
    draw.text((118, 925), "TEAM SHEEPDOG", font=font(25, bold=True), fill=MUTED)
    draw.text((118, 965), "appleboycat  |  du-du-lu", font=font(23), fill=MUTED)
    draw.text((118, 1005), "github.com/appleboycat/LexiPilot", font=font(22), fill=MUTED)
    return image


def slide_2() -> Image.Image:
    image, draw = slide_base(2, "Static vocabulary tools miss the learner", "The Problem")
    cards = [
        ("Static lists", "The same sequence is shown regardless of due dates, repeated mistakes, or available time.", RED),
        ("Fragmented memory", "Review history exists, but the learner still decides manually what matters today.", YELLOW),
        ("Weak adaptation", "A wrong answer rarely changes the next activity or generated practice.", PURPLE),
        ("Unsafe LLM control", "Unrestricted model writes can fabricate words, corrupt progress, or expose private context.", CYAN),
    ]
    boxes = [(90, 235, 910, 545), (1010, 235, 1830, 545), (90, 600, 910, 910), (1010, 600, 1830, 910)]
    for data, box in zip(cards, boxes):
        card(draw, box, *data)
    return image


def slide_3() -> Image.Image:
    image, draw = slide_base(3, "A hybrid Agent: adaptive decisions, controlled actions", "The Solution")
    steps = [
        ("1", "Natural-language goal", RED),
        ("2", "Read-only Tool Calling", CYAN),
        ("3", "Validated JSON plan", CYAN),
        ("4", "Deterministic session", GREEN),
        ("5", "Persistent feedback", PURPLE),
    ]
    x = 90
    for index, (number, label, color) in enumerate(steps):
        box = (x, 280, x + 300, 470)
        panel(draw, box, outline=color)
        draw.text((x + 25, 305), number, font=font(44, bold=True), fill=color)
        draw_wrapped(draw, (x + 25, 375), label, font(25, bold=True), WHITE, 245)
        if index < len(steps) - 1:
            draw.line((x + 300, 375, x + 345, 375), fill=MUTED, width=4)
            draw.polygon([(x + 345, 375), (x + 329, 365), (x + 329, 385)], fill=MUTED)
        x += 350
    card(
        draw,
        (90, 565, 910, 900),
        "Model responsibility",
        "Inspect learner state, choose read-only tools, and propose a plan grounded in returned evidence.",
        CYAN,
    )
    card(
        draw,
        (1010, 565, 1830, 900),
        "Controller responsibility",
        "Present cards, validate input, record explicit answers, apply spaced repetition, and finalize exactly once.",
        GREEN,
    )
    return image


def slide_4() -> Image.Image:
    image, draw = slide_base(4, "Hybrid Agent architecture and trust boundaries", "Architecture")
    panel(draw, (70, 220, 1850, 730), fill="#F8FAFB", outline=LINE)
    architecture = Image.open(ARCH_PNG).convert("RGB")
    architecture.thumbnail((1700, 455), Image.Resampling.LANCZOS)
    image.paste(architecture, (110 + (1700 - architecture.width) // 2, 245 + (455 - architecture.height) // 2))
    card(draw, (90, 780, 610, 950), "Read-only planning", "The model cannot record answers or save progress.", CYAN)
    card(draw, (700, 780, 1220, 950), "Validated execution", "Unknown words and malformed plans are rejected.", GREEN)
    card(draw, (1310, 780, 1830, 950), "Minimum disclosure", "Only task-relevant context crosses the endpoint boundary.", RED)
    return image


def slide_5() -> Image.Image:
    image, draw = slide_base(5, "Safety is enforced in code, not promised by the model", "Safe and Reliable Agent Design")
    rows = [
        ("Read-only planning tools", "No record_answer or save_session_summary schema is exposed.", CYAN),
        ("Strict plan validation", "Exact schema, real tool evidence, known words, and practical limits.", CYAN),
        ("Deterministic writes", "Only explicit y/n answers update spaced repetition.", GREEN),
        ("Failure containment", "Timeouts, fake Tool Calling, bad JSON, and unknown tools fall back.", YELLOW),
        ("Privacy-safe reporting", "No keys, private URLs, prompts, responses, or full profiles.", RED),
    ]
    y = 235
    for heading, body, accent in rows:
        panel(draw, (90, y, 1830, y + 125), outline=accent, fill=PANEL_ALT)
        draw.text((125, y + 28), heading, font=font(28, bold=True), fill=accent)
        draw_wrapped(draw, (650, y + 28), body, font(26), WHITE, 1125, line_gap=7)
        y += 145
    return image


def slide_6() -> Image.Image:
    image, draw = slide_base(6, "Qwen3-8B on a dedicated Radeon inference endpoint", "AMD Radeon Deployment")
    stack = [
        ("LexiPilot client", "OpenAI-compatible non-streaming requests", CYAN),
        ("vLLM server", "Hermes tool parser and automatic Tool Calling", GREEN),
        ("Qwen3-8B", "Planning and bilingual practice generation", PURPLE),
        ("ROCm + AMD Radeon", "Dedicated Cloud GPU instance", RED),
    ]
    y = 245
    for heading, detail, accent in stack:
        panel(draw, (90, y, 930, y + 135), outline=accent)
        draw.text((125, y + 26), heading, font=font(29, bold=True), fill=accent)
        draw.text((125, y + 78), detail, font=font(23), fill=WHITE)
        y += 155
    panel(draw, (1030, 245, 1830, 865), outline=RED)
    draw.text((1080, 290), "Deployment Evidence", font=font(34, bold=True), fill=RED)
    evidence = [
        "GPU model: Capture during final Radeon demo",
        "ROCm version: Capture during final Radeon demo",
        "vLLM version: Capture during final Radeon demo",
        "GPU activity: Capture during final Radeon demo",
        "Endpoint Tool Calling: verified by script",
    ]
    bullet_list(draw, evidence, (1080, 380, 1775, 820), accent=YELLOW, text_size=25, gap=20)
    draw.text((90, 900), "Learner state stays in the controlled LexiPilot environment; minimum task context is sent to the dedicated endpoint.", font=font(23), fill=MUTED)
    return image


def slide_7() -> Image.Image:
    image, draw = slide_base(7, "Measured thinking-mode comparison", "Evaluation")
    draw.text((90, 215), "1 warm-up + 5 measured requests per mode and workload | temperature 0 | max tokens 700", font=font(23), fill=MUTED)
    columns = [90, 750, 1190, 1570, 1830]
    headers = ["Workload / Metric", "Thinking enabled", "Thinking disabled", "Observed result"]
    for index, header in enumerate(headers):
        x1, x2 = columns[index], columns[index + 1]
        draw.rectangle((x1, 280, x2, 350), fill="#29343E", outline=LINE)
        draw.text((x1 + 18, 300), header, font=font(21, bold=True), fill=WHITE)
    rows = [
        ("Planning median latency", "13.7371 s", "13.7773 s", "No clear improvement"),
        ("Planning Tool Calling", "100%", "100%", "Same reliability"),
        ("Generation median latency", "7.0341 s", "6.5433 s", "6.98% lower median"),
        ("Generation validation", "100%", "100%", "Same reliability"),
        ("Generation completion tokens/s", "17.9127", "19.2563", "7.50% higher"),
    ]
    y = 350
    for row_index, row in enumerate(rows):
        fill = PANEL if row_index % 2 == 0 else PANEL_ALT
        for index, value in enumerate(row):
            x1, x2 = columns[index], columns[index + 1]
            draw.rectangle((x1, y, x2, y + 92), fill=fill, outline=LINE)
            color = CYAN if index in {1, 2} else GREEN if index == 3 and "lower" in value or "higher" in value else WHITE
            draw.text((x1 + 18, y + 30), value, font=font(22, bold=index > 0), fill=color)
        y += 92
    panel(draw, (90, 845, 1830, 955), outline=YELLOW)
    draw.text((120, 870), "Measurement scope", font=font(24, bold=True), fill=YELLOW)
    draw.text((420, 872), "Client-observed end-to-end values include network, endpoint, scheduling, and serving overhead - not raw GPU throughput.", font=font(21), fill=WHITE)
    return image


def slide_8() -> Image.Image:
    image, draw = slide_base(8, "A reproducible Agent pattern for private adaptive learning", "Impact and Roadmap")
    card(draw, (90, 235, 900, 560), "Submission impact", "A useful learning workflow with real persistent memory, safe model autonomy, deterministic writes, and Radeon-hosted inference.", GREEN)
    card(draw, (1020, 235, 1830, 560), "Reproducible by judges", "A sanitized 40-word index, synthetic profile, model-free smoke test, endpoint verifier, and benchmark script.", CYAN)
    panel(draw, (90, 625, 1830, 920), outline=RED)
    draw.text((130, 665), "Roadmap", font=font(34, bold=True), fill=RED)
    bullet_list(
        draw,
        [
            "Localhost AMD Radeon workstation deployment",
            "Optional GUI after the CLI trust model remains stable",
            "Multilingual practice beyond English and Chinese",
            "Richer long-term learner analytics",
        ],
        (130, 735, 1770, 900),
        accent=PURPLE,
        text_size=25,
        gap=10,
    )
    draw.text((1270, 958), "github.com/appleboycat/LexiPilot", font=font(24, bold=True), fill=CYAN)
    return image


SLIDE_BUILDERS = [slide_1, slide_2, slide_3, slide_4, slide_5, slide_6, slide_7, slide_8]


def build_slide_images(directory: Path) -> list[Path]:
    paths: list[Path] = []
    for index, builder in enumerate(SLIDE_BUILDERS, start=1):
        path = directory / f"slide_{index:02d}.png"
        builder().save(path, "PNG", optimize=True)
        paths.append(path)
    return paths


def build_presentation(slides: list[Path], destination: Path) -> None:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    presentation.core_properties.title = "LexiPilot Presentation"
    presentation.core_properties.author = "sheepdog - appleboycat, du-du-lu"
    presentation.core_properties.subject = "AMD Radeon Hackathon 2026 Track 2"
    blank = presentation.slide_layouts[6]
    for image_path in slides:
        slide = presentation.slides.add_slide(blank)
        slide.shapes.add_picture(str(image_path), 0, 0, width=presentation.slide_width, height=presentation.slide_height)
    temporary = destination.with_name(f".{destination.name}.tmp")
    presentation.save(temporary)
    atomic_replace(temporary, destination)


def build_presentation_pdf(slides: list[Path], destination: Path) -> None:
    from reportlab.pdfgen import canvas

    temporary = destination.with_name(f".{destination.name}.tmp")
    page_size = landscape((7.5 * inch, 13.333333 * inch))
    pdf = canvas.Canvas(str(temporary), pagesize=page_size)
    for image_path in slides:
        pdf.drawImage(str(image_path), 0, 0, width=page_size[0], height=page_size[1])
        pdf.showPage()
    pdf.setTitle("LexiPilot Presentation")
    pdf.setAuthor("sheepdog - appleboycat, du-du-lu")
    pdf.save()
    atomic_replace(temporary, destination)


def inline_markup(text: str) -> str:
    escaped = html.escape(text)
    escaped = re.sub(r"`([^`]+)`", r'<font name="Courier">\1</font>', escaped)
    escaped = re.sub(r"\*\*([^*]+)\*\*", r"<b>\1</b>", escaped)
    escaped = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", escaped)
    return escaped


def spec_styles() -> dict[str, ParagraphStyle]:
    base = getSampleStyleSheet()
    return {
        "title": ParagraphStyle(
            "SpecTitle", parent=base["Title"], fontName="Helvetica-Bold",
            fontSize=23, leading=27, textColor=colors.HexColor("#20252B"),
            alignment=TA_CENTER, spaceAfter=18,
        ),
        "h1": ParagraphStyle(
            "SpecH1", parent=base["Heading1"], fontName="Helvetica-Bold",
            fontSize=14, leading=17, textColor=colors.HexColor("#D52B1E"),
            spaceBefore=12, spaceAfter=7,
        ),
        "h2": ParagraphStyle(
            "SpecH2", parent=base["Heading2"], fontName="Helvetica-Bold",
            fontSize=10.5, leading=13, textColor=colors.HexColor("#147D95"),
            spaceBefore=8, spaceAfter=4,
        ),
        "body": ParagraphStyle(
            "SpecBody", parent=base["BodyText"], fontName="Helvetica",
            fontSize=8.5, leading=11, textColor=colors.HexColor("#20252B"),
            alignment=TA_LEFT, spaceAfter=5,
        ),
        "bullet": ParagraphStyle(
            "SpecBullet", parent=base["BodyText"], fontName="Helvetica",
            fontSize=8.5, leading=11, leftIndent=14, firstLineIndent=-7,
            bulletIndent=6, spaceAfter=2,
        ),
        "quote": ParagraphStyle(
            "SpecQuote", parent=base["BodyText"], fontName="Helvetica-Oblique",
            fontSize=8.5, leading=11, leftIndent=18, rightIndent=18,
            textColor=colors.HexColor("#4B5563"), spaceAfter=6,
        ),
        "code": ParagraphStyle(
            "SpecCode", parent=base["Code"], fontName="Courier", fontSize=6.8,
            leading=8.5, leftIndent=8, rightIndent=8, backColor=colors.HexColor("#F3F5F7"),
            borderColor=colors.HexColor("#D7DDE2"), borderWidth=0.5,
            borderPadding=6, spaceAfter=7,
        ),
        "small": ParagraphStyle(
            "SpecSmall", parent=base["BodyText"], fontName="Helvetica",
            fontSize=7.2, leading=9, textColor=colors.HexColor("#4B5563"),
        ),
    }


def table_flowable(rows: list[list[str]], available_width: float, styles: dict[str, ParagraphStyle]) -> Table:
    column_count = max(len(row) for row in rows)
    normalized = [row + [""] * (column_count - len(row)) for row in rows]
    data = [[Paragraph(inline_markup(cell), styles["small"]) for cell in row] for row in normalized]
    table = Table(data, colWidths=[available_width / column_count] * column_count, repeatRows=1)
    table.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#25313A")),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("BACKGROUND", (0, 1), (-1, -1), colors.HexColor("#F7F8FA")),
                ("GRID", (0, 0), (-1, -1), 0.4, colors.HexColor("#C5CDD4")),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ("LEFTPADDING", (0, 0), (-1, -1), 4),
                ("RIGHTPADDING", (0, 0), (-1, -1), 4),
                ("TOPPADDING", (0, 0), (-1, -1), 3),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 3),
            ]
        )
    )
    return table


def markdown_story(markdown: str, available_width: float) -> list[Any]:
    styles = spec_styles()
    story: list[Any] = []
    paragraph_lines: list[str] = []
    code_lines: list[str] = []
    table_rows: list[list[str]] = []
    in_code = False

    def flush_paragraph() -> None:
        if paragraph_lines:
            story.append(Paragraph(inline_markup(" ".join(paragraph_lines)), styles["body"]))
            paragraph_lines.clear()

    def flush_table() -> None:
        if table_rows:
            cleaned = [
                row for row in table_rows
                if not all(re.fullmatch(r":?-{3,}:?", cell.strip()) for cell in row)
            ]
            if cleaned:
                story.append(table_flowable(cleaned, available_width, styles))
                story.append(Spacer(1, 6))
            table_rows.clear()

    for raw in markdown.splitlines():
        line = raw.rstrip()
        if line.startswith("```"):
            flush_paragraph()
            flush_table()
            if in_code:
                story.append(Preformatted("\n".join(code_lines), styles["code"]))
                code_lines.clear()
                in_code = False
            else:
                in_code = True
            continue
        if in_code:
            code_lines.append(line)
            continue
        if line.startswith("|") and line.endswith("|"):
            flush_paragraph()
            table_rows.append([cell.strip() for cell in line.strip("|").split("|")])
            continue
        flush_table()
        if not line.strip():
            flush_paragraph()
            continue
        if line.startswith("# "):
            flush_paragraph()
            story.append(Paragraph(inline_markup(line[2:]), styles["title"]))
        elif line.startswith("## "):
            flush_paragraph()
            story.append(Paragraph(inline_markup(line[3:]), styles["h1"]))
        elif line.startswith("### "):
            flush_paragraph()
            story.append(Paragraph(inline_markup(line[4:]), styles["h2"]))
        elif line.startswith("![") and ARCH_PNG.name in line:
            flush_paragraph()
            architecture = RLImage(str(ARCH_PNG))
            architecture.drawWidth = available_width
            architecture.drawHeight = available_width * 0.24
            story.append(KeepTogether([architecture, Spacer(1, 7)]))
        elif line.startswith("- "):
            flush_paragraph()
            story.append(Paragraph(inline_markup(line[2:]), styles["bullet"], bulletText="-"))
        elif re.match(r"^\d+\. ", line):
            flush_paragraph()
            number, text = line.split(". ", 1)
            story.append(Paragraph(inline_markup(text), styles["bullet"], bulletText=f"{number}."))
        elif line.startswith("> "):
            flush_paragraph()
            story.append(Paragraph(inline_markup(line[2:]), styles["quote"]))
        elif line == "---":
            flush_paragraph()
            story.append(Spacer(1, 5))
        else:
            paragraph_lines.append(line.strip())
    flush_paragraph()
    flush_table()
    return story


def build_spec_pdf(destination: Path) -> None:
    temporary = destination.with_name(f".{destination.name}.tmp")
    page_width, _ = A4
    left_margin = right_margin = 0.55 * inch
    document = SimpleDocTemplate(
        str(temporary),
        pagesize=A4,
        leftMargin=left_margin,
        rightMargin=right_margin,
        topMargin=0.55 * inch,
        bottomMargin=0.55 * inch,
        title="LexiPilot Project Specification",
        author="sheepdog - appleboycat, du-du-lu",
        subject="AMD Radeon Hackathon 2026 Track 2",
    )
    story = markdown_story(SPEC_MD.read_text(encoding="utf-8"), page_width - left_margin - right_margin)

    def page_decor(canvas: Any, doc: Any) -> None:
        canvas.saveState()
        canvas.setStrokeColor(colors.HexColor("#D52B1E"))
        canvas.setLineWidth(2)
        canvas.line(doc.leftMargin, A4[1] - 24, A4[0] - doc.rightMargin, A4[1] - 24)
        canvas.setFont("Helvetica", 7)
        canvas.setFillColor(colors.HexColor("#5A6570"))
        canvas.drawString(doc.leftMargin, 18, "LexiPilot | AMD Radeon Hackathon 2026 | Track 2")
        canvas.drawRightString(A4[0] - doc.rightMargin, 18, f"Page {doc.page}")
        canvas.restoreState()

    document.build(story, onFirstPage=page_decor, onLaterPages=page_decor)
    atomic_replace(temporary, destination)


def build_all() -> None:
    if not ARCH_PNG.exists():
        raise SystemExit(f"Missing architecture image: {ARCH_PNG}")
    with tempfile.TemporaryDirectory(prefix="lexipilot_submission_") as temp_name:
        temp = Path(temp_name)
        slides = build_slide_images(temp)
        build_presentation(slides, PRESENTATION)
        build_presentation_pdf(slides, PRESENTATION_PDF)
    build_spec_pdf(SPEC_PDF)

    OFFICIAL.mkdir(parents=True, exist_ok=True)
    for source, name in (
        (SPEC_PDF, "project_specification.pdf"),
        (PRESENTATION_PDF, "presentation.pdf"),
    ):
        temporary = OFFICIAL / f".{name}.tmp"
        shutil.copy2(source, temporary)
        atomic_replace(temporary, OFFICIAL / name)

    for path in (SPEC_PDF, PRESENTATION, PRESENTATION_PDF, OFFICIAL / "project_specification.pdf", OFFICIAL / "presentation.pdf"):
        print(f"PASS built {path.relative_to(ROOT)}")


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.parse_args()
    build_all()
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
